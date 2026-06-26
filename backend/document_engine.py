"""
document_engine.py — Quantum MCAGI
=====================================
Unified document handling engine.
Merges: document_parser.py + document_ingester.py

Supports:
  Local files: .txt .md .py .pdf .docx .xlsx .csv .json .html .xml .pptx
  URLs:        Any webpage, direct file link
  Cloud:       Google Drive, Dropbox, OneDrive shared links
  Images:      .jpg .png (OCR if available)
"""
import os
import io
import csv
import json
import re

def parse_document(file_path=None, file_bytes=None, filename=None):
    if file_path:
        filename = filename or os.path.basename(file_path)
        with open(file_path, 'rb') as f:
            file_bytes = f.read()
    
    if not filename or not file_bytes:
        return {"error": "No file provided", "text": ""}
    
    ext = os.path.splitext(filename)[1].lower()
    
    parsers = {
        '.txt': _parse_txt,
        '.md': _parse_txt,
        '.py': _parse_txt,
        '.js': _parse_txt,
        '.ts': _parse_txt,
        '.tsx': _parse_txt,
        '.jsx': _parse_txt,
        '.json': _parse_json,
        '.csv': _parse_csv,
        '.docx': _parse_docx,
        '.doc': _parse_doc,
        '.pdf': _parse_pdf,
        '.xlsx': _parse_xlsx,
        '.xls': _parse_xlsx,
        '.ods': _parse_ods,
        '.html': _parse_html,
        '.htm': _parse_html,
        '.xml': _parse_xml,
        '.rtf': _parse_rtf,
        '.yaml': _parse_txt,
        '.yml': _parse_txt,
        '.toml': _parse_txt,
        '.ini': _parse_txt,
        '.cfg': _parse_txt,
        '.log': _parse_txt,
        '.sh': _parse_txt,
        '.bash': _parse_txt,
        '.bat': _parse_txt,
        '.sql': _parse_txt,
        '.r': _parse_txt,
        '.java': _parse_txt,
        '.c': _parse_txt,
        '.cpp': _parse_txt,
        '.h': _parse_txt,
        '.rb': _parse_txt,
        '.go': _parse_txt,
        '.rs': _parse_txt,
        '.swift': _parse_txt,
        '.kt': _parse_txt,
        '.lua': _parse_txt,
        '.php': _parse_txt,
        '.pl': _parse_txt,
        '.css': _parse_txt,
        '.scss': _parse_txt,
        '.sass': _parse_txt,
        '.less': _parse_txt,
    }
    
    parser = parsers.get(ext)
    if not parser:
        try:
            text = file_bytes.decode('utf-8', errors='replace')
            return {"text": text, "format": "unknown_text", "chars": len(text)}
        except Exception:
            return {"error": f"Unsupported format: {ext}", "text": ""}
    
    try:
        result = parser(file_bytes, filename)
        result['format'] = ext
        return result
    except Exception as e:
        return {"error": f"Parse error ({ext}): {str(e)}", "text": ""}


def _parse_txt(data, filename=""):
    text = data.decode('utf-8', errors='replace')
    return {"text": text, "chars": len(text)}


def _parse_json(data, filename=""):
    text = data.decode('utf-8', errors='replace')
    try:
        obj = json.loads(text)
        flat = _flatten_json(obj)
        return {"text": flat, "chars": len(flat)}
    except Exception:
        return {"text": text, "chars": len(text)}


def _flatten_json(obj, prefix=""):
    parts = []
    if isinstance(obj, dict):
        for k, v in obj.items():
            if isinstance(v, str):
                parts.append(f"{k}: {v}")
            elif isinstance(v, (int, float, bool)):
                parts.append(f"{k}: {v}")
            else:
                parts.append(_flatten_json(v, f"{prefix}{k}."))
    elif isinstance(obj, list):
        for item in obj:
            if isinstance(item, str):
                parts.append(item)
            else:
                parts.append(_flatten_json(item, prefix))
    else:
        parts.append(str(obj))
    return "\n".join(parts)


def _parse_csv(data, filename=""):
    text = data.decode('utf-8', errors='replace')
    reader = csv.reader(io.StringIO(text))
    rows = []
    for row in reader:
        rows.append(" ".join(row))
    result = "\n".join(rows)
    return {"text": result, "chars": len(result), "rows": len(rows)}


def _parse_docx(data, filename=""):
    from docx import Document
    doc = Document(io.BytesIO(data))
    paragraphs = []
    for para in doc.paragraphs:
        t = para.text.strip()
        if t:
            paragraphs.append(t)
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                paragraphs.append(" | ".join(cells))
    text = "\n\n".join(paragraphs)
    return {"text": text, "chars": len(text), "paragraphs": len(paragraphs)}


def _parse_doc(data, filename=""):
    try:
        text = data.decode('utf-8', errors='replace')
        text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f]', ' ', text)
        text = re.sub(r'\s+', ' ', text).strip()
        return {"text": text, "chars": len(text)}
    except Exception:
        return {"error": "Legacy .doc format — save as .docx for best results", "text": ""}


def _parse_pdf(data, filename=""):
    from PyPDF2 import PdfReader
    reader = PdfReader(io.BytesIO(data))
    pages = []
    for page in reader.pages:
        t = page.extract_text()
        if t:
            pages.append(t.strip())
    text = "\n\n".join(pages)
    return {"text": text, "chars": len(text), "pages": len(pages)}


def _parse_xlsx(data, filename=""):
    from openpyxl import load_workbook
    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    rows = []
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        rows.append(f"Sheet: {sheet}")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(cells):
                rows.append(" | ".join(cells))
    text = "\n".join(rows)
    return {"text": text, "chars": len(text), "rows": len(rows)}


def _parse_ods(data, filename=""):
    from odf.opendocument import load as odf_load
    from odf.table import Table, TableRow, TableCell
    from odf.text import P
    doc = odf_load(io.BytesIO(data))
    rows = []
    for table in doc.getElementsByType(Table):
        rows.append(f"Sheet: {table.getAttribute('name')}")
        for row in table.getElementsByType(TableRow):
            cells = []
            for cell in row.getElementsByType(TableCell):
                text_parts = []
                for p in cell.getElementsByType(P):
                    t = ""
                    for child in p.childNodes:
                        t += str(child)
                    text_parts.append(t)
                cells.append(" ".join(text_parts))
            if any(cells):
                rows.append(" | ".join(cells))
    text = "\n".join(rows)
    return {"text": text, "chars": len(text), "rows": len(rows)}


def _parse_html(data, filename=""):
    text = data.decode('utf-8', errors='replace')
    text = re.sub(r'<script[^>]*>.*?</script>', '', text, flags=re.DOTALL)
    text = re.sub(r'<style[^>]*>.*?</style>', '', text, flags=re.DOTALL)
    text = re.sub(r'<[^>]+>', ' ', text)
    text = re.sub(r'\s+', ' ', text).strip()
    return {"text": text, "chars": len(text)}


def _parse_xml(data, filename=""):
    text = data.decode('utf-8', errors='replace')
    text = re.sub(r'<[^>]+>', ' ', text)
    text = re.sub(r'\s+', ' ', text).strip()
    return {"text": text, "chars": len(text)}


def _parse_rtf(data, filename=""):
    text = data.decode('utf-8', errors='replace')
    text = re.sub(r'\\[a-z]+\d*\s?', '', text)
    text = re.sub(r'[{}]', '', text)
    text = re.sub(r'\s+', ' ', text).strip()
    return {"text": text, "chars": len(text)}


SUPPORTED_FORMATS = [
    '.txt', '.md', '.py', '.js', '.ts', '.tsx', '.json', '.csv',
    '.docx', '.doc', '.pdf', '.xlsx', '.xls', '.ods',
    '.html', '.htm', '.xml', '.rtf', '.yaml', '.yml', '.toml',
    '.sh', '.sql', '.java', '.c', '.cpp', '.go', '.rs', '.rb',
    '.php', '.css', '.scss', '.log', '.ini', '.cfg',
]

"""
Quantum MCAGI — Document Ingester
==================================
No document discrimination. Every format welcome.

Supports:
  Local files:  .txt .md .py .pdf .docx .xlsx .csv .pptx .json .html .xml
  URLs:         Any webpage, direct file link
  Google Drive: Shared links (public)
  Dropbox:      Shared links
  OneDrive:     Shared links
  Images:       .jpg .png .gif .bmp (text extraction via OCR if available)
"""

import os
import re
import json
import requests
import urllib.parse
from pathlib import Path
from typing import Optional

# ── Format handlers ──────────────────────────────────────────────────────────

def extract_txt(filepath):
    with open(filepath, 'r', errors='ignore') as f:
        return f.read()

def extract_pdf(filepath):
    try:
        import PyPDF2
        text = []
        with open(filepath, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                text.append(page.extract_text() or '')
        return "\n".join(text)
    except Exception as e:
        return f"[PDF ERROR: {e}]"

def extract_docx(filepath):
    try:
        import docx
        doc = docx.Document(filepath)
        return "\n".join(p.text for p in doc.paragraphs)
    except Exception as e:
        return f"[DOCX ERROR: {e}]"

def extract_xlsx(filepath):
    try:
        import openpyxl
        wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
        text = []
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            text.append(f"[Sheet: {sheet}]")
            for row in ws.iter_rows(values_only=True):
                row_text = ' '.join(str(c) for c in row if c is not None)
                if row_text.strip():
                    text.append(row_text)
        return "\n".join(text)
    except Exception as e:
        return f"[XLSX ERROR: {e}]"

def extract_pptx(filepath):
    try:
        from pptx import Presentation
        prs = Presentation(filepath)
        text = []
        for i, slide in enumerate(prs.slides):
            text.append(f"[Slide {i+1}]")
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    text.append(shape.text)
        return "\n".join(text)
    except Exception as e:
        return f"[PPTX ERROR: {e}]"

def extract_csv(filepath):
    import csv
    text = []
    try:
        with open(filepath, 'r', errors='ignore') as f:
            reader = csv.reader(f)
            for row in reader:
                text.append(' '.join(row))
        return "\n".join(text)
    except Exception as e:
        return f"[CSV ERROR: {e}]"

def extract_image(filepath):
    try:
        import pytesseract
        from PIL import Image
        img = Image.open(filepath)
        return pytesseract.image_to_string(img)
    except ImportError:
        return "[IMAGE: OCR not available — install pytesseract]"
    except Exception as e:
        return f"[IMAGE ERROR: {e}]"

def extract_html(filepath=None, html_text=None):
    try:
        from bs4 import BeautifulSoup
        if filepath:
            with open(filepath, 'r', errors='ignore') as f:
                html_text = f.read()
        soup = BeautifulSoup(html_text, 'html.parser')
        for tag in soup(['script', 'style', 'nav', 'footer']):
            tag.decompose()
        return soup.get_text(separator="\n", strip=True)
    except Exception as e:
        return f"[HTML ERROR: {e}]"

def extract_json(filepath):
    try:
        with open(filepath, 'r', errors='ignore') as f:
            data = json.load(f)
        return json.dumps(data, indent=2)
    except Exception as e:
        return f"[JSON ERROR: {e}]"

# ── Extension router ──────────────────────────────────────────────────────────

EXTRACTORS = {
    '.txt': extract_txt,
    '.md': extract_txt,
    '.py': extract_txt,
    '.js': extract_txt,
    '.xml': extract_txt,
    '.rst': extract_txt,
    '.pdf': extract_pdf,
    '.docx': extract_docx,
    '.doc': extract_docx,
    '.xlsx': extract_xlsx,
    '.xls': extract_xlsx,
    '.pptx': extract_pptx,
    '.ppt': extract_pptx,
    '.csv': extract_csv,
    '.tsv': extract_csv,
    '.html': extract_html,
    '.htm': extract_html,
    '.json': extract_json,
    '.jpg': extract_image,
    '.jpeg': extract_image,
    '.png': extract_image,
    '.gif': extract_image,
    '.bmp': extract_image,
}

def extract_from_file(filepath):
    """Extract text from any local file."""
    path = Path(filepath)
    if not path.exists():
        return None, f"File not found: {filepath}"
    ext = path.suffix.lower()
    extractor = EXTRACTORS.get(ext, extract_txt)
    try:
        text = extractor(str(filepath))
        return text, f"extracted {len(text.split()):,} words from {path.name}"
    except Exception as e:
        return None, str(e)

# ── URL & Cloud handlers ──────────────────────────────────────────────────────

HEADERS = {'User-Agent': 'Mozilla/5.0 (X11; Linux x86_64) AppleWebKit/537.36'}

def resolve_drive_url(url):
    """Convert Google Drive share link to direct download."""
    match = re.search(r'/file/d/([^/]+)', url)
    if match:
        fid = match.group(1)
        return f"https://drive.google.com/uc?export=download&id={fid}"
    match = re.search(r'id=([^&]+)', url)
    if match:
        return f"https://drive.google.com/uc?export=download&id={match.group(1)}"
    return url

def resolve_dropbox_url(url):
    """Convert Dropbox share link to direct download."""
    return url.replace('www.dropbox.com', 'dl.dropboxusercontent.com').replace('?dl=0', '').replace('?dl=1', '')

def resolve_onedrive_url(url):
    """Convert OneDrive share link to direct download."""
    if '1drv.ms' in url or 'onedrive.live.com' in url:
        encoded = urllib.parse.quote(url, safe='')
        return f"https://api.onedrive.com/v1.0/shares/u!{encoded}/root/content"
    return url

def resolve_url(url):
    """Resolve cloud share links to direct download URLs."""
    if 'drive.google.com' in url:
        return resolve_drive_url(url), 'gdrive'
    elif 'dropbox.com' in url:
        return resolve_dropbox_url(url), 'dropbox'
    elif '1drv.ms' in url or 'onedrive.live.com' in url:
        return resolve_onedrive_url(url), 'onedrive'
    return url, 'direct'

def fetch_url(url, max_chars=100000):
    """Fetch text content from any URL. Smart extraction for known sources."""
    resolved_url, source_type = resolve_url(url)

    try:
        r = requests.get(resolved_url, headers=HEADERS, timeout=30, allow_redirects=True)
        if r.status_code != 200:
            return None, f"HTTP {r.status_code}"

        content_type = r.headers.get('content-type', '').lower()

        if 'pdf' in content_type or url.lower().endswith('.pdf'):
            import tempfile
            with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as tmp:
                tmp.write(r.content)
                tmp_path = tmp.name
            text, status = extract_from_file(tmp_path)
            os.unlink(tmp_path)
            return text, status

        elif 'word' in content_type or 'docx' in content_type or url.lower().endswith('.docx'):
            import tempfile
            with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
                tmp.write(r.content)
                tmp_path = tmp.name
            text, status = extract_from_file(tmp_path)
            os.unlink(tmp_path)
            return text, status

        elif 'spreadsheet' in content_type or url.lower().endswith('.xlsx'):
            import tempfile
            with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
                tmp.write(r.content)
                tmp_path = tmp.name
            text, status = extract_from_file(tmp_path)
            os.unlink(tmp_path)
            return text, status

        elif 'html' in content_type:
            text = _smart_extract_html(url, r.text)
            return text[:max_chars], f"extracted {len(text.split()):,} words from webpage"

        elif 'json' in content_type:
            return r.text[:max_chars], f"extracted JSON ({len(r.text.split()):,} words)"

        else:
            text = r.text[:max_chars]
            return text, f"extracted {len(text.split()):,} words"

    except Exception as e:
        return None, str(e)


def _smart_extract_html(url, html_text):
    """
    Smart content extraction for known sources.
    Extracts article body text, strips navigation/sidebars/footers.
    Falls back to generic extraction for unknown sources.
    """
    try:
        from bs4 import BeautifulSoup
    except ImportError:
        return extract_html(html_text=html_text)

    soup = BeautifulSoup(html_text, 'html.parser')

    # Wikipedia
    if 'wikipedia.org' in url:
        content = soup.find('div', {'id': 'mw-content-text'})
        if content:
            for tag in content.find_all(['table', 'sup', 'span'], class_=lambda c: c and any(
                x in str(c) for x in ['navbox', 'sidebar', 'reflist', 'mw-editsection', 'noprint', 'metadata'])):
                tag.decompose()
            for tag in content.find_all('div', class_=lambda c: c and any(
                x in str(c) for x in ['reflist', 'navbox', 'sistersitebox', 'noprint'])):
                tag.decompose()
            return content.get_text(separator="\n", strip=True)

    # Stanford Encyclopedia of Philosophy
    if 'plato.stanford.edu' in url:
        content = soup.find('div', {'id': 'aueditable'}) or soup.find('div', {'id': 'article-content'})
        if content:
            for tag in content.find_all(['nav', 'footer']):
                tag.decompose()
            return content.get_text(separator="\n", strip=True)

    # World History Encyclopedia
    if 'worldhistory.org' in url:
        content = soup.find('div', class_='article-content-main') or soup.find('article') or soup.find('div', {'id': 'content'})
        if content:
            for tag in content.find_all(['nav', 'footer', 'aside']):
                tag.decompose()
            return content.get_text(separator="\n", strip=True)

    # arXiv
    if 'arxiv.org' in url:
        abstract = soup.find('blockquote', class_='abstract')
        title = soup.find('h1', class_='title')
        parts = []
        if title:
            parts.append(title.get_text(strip=True))
        if abstract:
            parts.append(abstract.get_text(strip=True))
        if parts:
            return "\n".join(parts)

    # Project Gutenberg
    if 'gutenberg.org' in url:
        content = soup.find('div', {'id': 'pg-machine-header'})
        if content:
            content.decompose()
        body = soup.find('body')
        if body:
            for tag in body.find_all(['table', 'pre']):
                if 'START OF' in tag.get_text() or 'END OF' in tag.get_text():
                    tag.decompose()
            return body.get_text(separator="\n", strip=True)

    # Generic fallback
    for tag in soup(['script', 'style', 'nav', 'footer', 'aside', 'header']):
        tag.decompose()
    main = soup.find('main') or soup.find('article') or soup.find('div', {'role': 'main'})
    if main:
        return main.get_text(separator="\n", strip=True)

    return soup.get_text(separator="\n", strip=True)


# ── Main interface ────────────────────────────────────────────────────────────

def ingest_document(source, max_chars=100000):
    """
    Universal document ingestion.
    source can be: local filepath, URL, Google Drive link, Dropbox link, OneDrive link
    Returns (text, status_message)
    """
    source = source.strip()

    if source.startswith('http://') or source.startswith('https://'):
        return fetch_url(source, max_chars)

    if os.path.exists(source):
        return extract_from_file(source)

    for base in ['/sdcard/Download/', '/sdcard/', os.path.expanduser('~/'), '']:
        full_path = os.path.join(base, source)
        if os.path.exists(full_path):
            return extract_from_file(full_path)

    return None, f"Cannot find: {source}"


def handle_ingest_command(cmd_parts, engine=None, memory=None):
    """
    Handle /ingest commands from chat.py

    /ingest URL                 — fetch and train from URL
    /ingest /path/to/file       — train from local file
    /ingest /sdcard/Download/x  — train from phone storage
    """
    if len(cmd_parts) < 2:
        return """  Ingest Commands:
  /ingest URL              — fetch webpage or file from URL
  /ingest FILEPATH         — ingest local file (any format)
  /ingest /sdcard/Download/FILE — ingest from phone Downloads

  Supported: .txt .pdf .docx .xlsx .pptx .csv .html .json .md .py .png .jpg
  Cloud:     Google Drive, Dropbox, OneDrive share links"""

    source = ' '.join(cmd_parts[1:])
    print(f"  Ingesting: {source[:60]}...")

    text, status = ingest_document(source)

    if not text:
        return f"  Failed: {status}"

    if engine:
        engine.learn_from_text(text)

        if memory:
            try:
                from quantum_language_engine import ConceptExtractor
                extractor = ConceptExtractor()
                extractor.update_corpus_stats(text)
                concepts = extractor.extract_concepts(text, max_concepts=10)
                for c in concepts:
                    name = c if isinstance(c, str) else c.get('concept', '')
                    if name and memory.concepts is not None and name not in memory.concepts:
                        memory.concepts[name] = {'count': 1, 'strength': 1.0}
                        memory.growth['total_concepts'] = memory.growth.get('total_concepts', 0) + 1
                memory.growth['total_insights'] = memory.growth.get('total_insights', 0) + 1
            except Exception:
                pass

        # Extract and persist SVO facts
        new_facts = 0
        try:
            fact_path = os.path.join(os.path.expanduser('~/.quantum-mcagi'), 'fact_store.json')
            try:
                with open(fact_path) as f:
                    fact_store = json.load(f)
            except Exception:
                fact_store = {}

            ASSERTION_VERBS = {'is', 'are', 'was', 'were', 'has', 'have', 'contains',
                               'includes', 'consists', 'means', 'refers', 'represents',
                               'located', 'found', 'known', 'called', 'defined'}
            SKIP_SUBJECTS = {'this book', 'the book', 'our company', 'the cover',
                             'the author', 'the publisher', 'this text', 'the text',
                             'this page', 'the page', 'we', 'it', 'this', 'that', 'he', 'she', 'they'}

            content = text[2000:]
            sentences = re.split(r'[.!?]+', content)
            for sentence in sentences:
                words = sentence.strip().split()
                if len(words) < 4 or len(words) > 40:
                    continue
                for i, word in enumerate(words):
                    if word.lower() in ASSERTION_VERBS and i >= 1:
                        subject = ' '.join(words[:i]).strip().lower()
                        verb = word.lower()
                        obj = ' '.join(words[i+1:]).strip().lower()
                        subj_words = subject.split()
                        if (subject and obj and 1 <= len(subj_words) <= 5
                                and subject not in SKIP_SUBJECTS
                                and not subject.startswith('this ')
                                and not subject.startswith('the following')
                                and len(obj) > 10
                                and not any(c.isdigit() for c in subject)):
                            if subject not in fact_store:
                                fact_store[subject] = []
                            triple = [verb, obj[:120]]
                            if triple not in fact_store[subject]:
                                fact_store[subject].append(triple)
                                new_facts += 1
                        break

            with open(fact_path, 'w') as f:
                json.dump(fact_store, f)
        except Exception:
            new_facts = -1

        word_count = len(text.split())

        # Detect word-list format (dictionary/vocab files) and add as concepts
        lines = text.strip().split('\n')
        is_word_list = sum(1 for l in lines[:200] if l.strip() and len(l.strip().split()) == 1) > 150
        concepts_added = 0
        if is_word_list and (memory or engine):
            all_words = [l.strip() for l in lines if l.strip().isalpha() and len(l.strip()) >= 3]
            concepts_added = ingest_words_as_knowledge(all_words, memory, engine)

        # Save TF-IDF so vocab persists across restarts
        tfidf_saved = ingest_save_tfidf(engine) if engine else False

        vocab_size = 0
        try:
            vocab_size = len(engine.tfidf.extractor.word_frequencies)
        except Exception:
            pass

        fact_msg = f" | Facts: +{new_facts:,}" if new_facts >= 0 else ""
        concept_msg = f" | Concepts: +{concepts_added:,}" if concepts_added > 0 else ""
        vocab_msg = f" | Vocab: {vocab_size:,}" if vocab_size > 0 else ""
        tfidf_msg = " | TF-IDF saved ✓" if tfidf_saved else ""
        return f"  ✓ {status} | Words: +{word_count:,} | States: {len(engine.markov.chain):,}{fact_msg}{concept_msg}{vocab_msg}{tfidf_msg}"

    return f"  ✓ {status} (no engine — text extracted only)"


if __name__ == "__main__":
    import sys
    if len(sys.argv) > 1:
        source = ' '.join(sys.argv[1:])
        text, status = ingest_document(source)
        print(f"Status: {status}")
        if text:
            print(f"Preview: {text[:500]}")
    else:
        print("Usage: python document_ingester.py URL_OR_FILEPATH")

def ingest_save_tfidf(engine):
    """Force-save TF-IDF state after ingest so vocab persists across restarts."""
    try:
        import pickle, os
        tfidf_path = os.path.join(os.path.expanduser('~/.quantum-mcagi'),
                                   'engine_state', 'tfidf_state.pkl')
        with open(tfidf_path, 'wb') as f:
            pickle.dump(engine.tfidf, f)
        return True
    except Exception as e:
        return False

def ingest_words_as_knowledge(words, memory, engine):
    """Add dictionary words directly as concepts and knowledge entries."""
    added = 0
    # Feed words into TF-IDF as a single corpus document
    text_block = ' '.join(words)
    if engine:
        engine.tfidf.learn(text_block)

    if not memory:
        return added

    for word in words:
        word = word.strip().lower()
        if not word or len(word) < 3 or not word.isalpha():
            continue
        if memory.concepts is not None and word not in memory.concepts:
            memory.concepts[word] = {
                'count': 1,
                'strength': 1.0,
                'source': 'oxford_dictionary',
                'relationships': {}
            }
            memory.growth['total_concepts'] = memory.growth.get('total_concepts', 0) + 1
            added += 1
    return added
