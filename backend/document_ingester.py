"""
Quantum MCAGI — Document Ingester
==================================
No document discrimination. Every format welcome.

Supports:
  Local files:  .txt .md .py .pdf .docx .xlsx .csv .pptx .json .html .xml
  URLs:         Any webpage, direct file link
  Google Drive: Shared links (public)
  Dropbox:      Shared links
  OneDrive:     Shared links
  Images:       .jpg .png .gif .bmp (text extraction via OCR if available)
"""

import os
import re
import json
import requests
import urllib.parse
from pathlib import Path
from typing import Optional

# ── Format handlers ──────────────────────────────────────────────────────────

def extract_txt(filepath):
    with open(filepath, 'r', errors='ignore') as f:
        return f.read()

def extract_pdf(filepath):
    try:
        import PyPDF2
        text = []
        with open(filepath, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                text.append(page.extract_text() or '')
        return '\n'.join(text)
    except Exception as e:
        return f"[PDF ERROR: {e}]"

def extract_docx(filepath):
    try:
        import docx
        doc = docx.Document(filepath)
        return '\n'.join(p.text for p in doc.paragraphs)
    except Exception as e:
        return f"[DOCX ERROR: {e}]"

def extract_xlsx(filepath):
    try:
        import openpyxl
        wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
        text = []
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            text.append(f"[Sheet: {sheet}]")
            for row in ws.iter_rows(values_only=True):
                row_text = ' '.join(str(c) for c in row if c is not None)
                if row_text.strip():
                    text.append(row_text)
        return '\n'.join(text)
    except Exception as e:
        return f"[XLSX ERROR: {e}]"

def extract_pptx(filepath):
    try:
        from pptx import Presentation
        prs = Presentation(filepath)
        text = []
        for i, slide in enumerate(prs.slides):
            text.append(f"[Slide {i+1}]")
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    text.append(shape.text)
        return '\n'.join(text)
    except Exception as e:
        return f"[PPTX ERROR: {e}]"

def extract_csv(filepath):
    import csv
    text = []
    try:
        with open(filepath, 'r', errors='ignore') as f:
            reader = csv.reader(f)
            for row in reader:
                text.append(' '.join(row))
        return '\n'.join(text)
    except Exception as e:
        return f"[CSV ERROR: {e}]"

def extract_image(filepath):
    try:
        import pytesseract
        from PIL import Image
        img = Image.open(filepath)
        return pytesseract.image_to_string(img)
    except ImportError:
        return "[IMAGE: OCR not available — install pytesseract]"
    except Exception as e:
        return f"[IMAGE ERROR: {e}]"

def extract_html(filepath=None, html_text=None):
    try:
        from bs4 import BeautifulSoup
        if filepath:
            with open(filepath, 'r', errors='ignore') as f:
                html_text = f.read()
        soup = BeautifulSoup(html_text, 'html.parser')
        for tag in soup(['script', 'style', 'nav', 'footer']):
            tag.decompose()
        return soup.get_text(separator='\n', strip=True)
    except Exception as e:
        return f"[HTML ERROR: {e}]"

def extract_json(filepath):
    try:
        with open(filepath, 'r', errors='ignore') as f:
            data = json.load(f)
        return json.dumps(data, indent=2)
    except Exception as e:
        return f"[JSON ERROR: {e}]"

# ── Extension router ──────────────────────────────────────────────────────────

EXTRACTORS = {
    '.txt': extract_txt,
    '.md': extract_txt,
    '.py': extract_txt,
    '.js': extract_txt,
    '.xml': extract_txt,
    '.rst': extract_txt,
    '.pdf': extract_pdf,
    '.docx': extract_docx,
    '.doc': extract_docx,
    '.xlsx': extract_xlsx,
    '.xls': extract_xlsx,
    '.pptx': extract_pptx,
    '.ppt': extract_pptx,
    '.csv': extract_csv,
    '.tsv': extract_csv,
    '.html': extract_html,
    '.htm': extract_html,
    '.json': extract_json,
    '.jpg': extract_image,
    '.jpeg': extract_image,
    '.png': extract_image,
    '.gif': extract_image,
    '.bmp': extract_image,
}

def extract_from_file(filepath):
    """Extract text from any local file."""
    path = Path(filepath)
    if not path.exists():
        return None, f"File not found: {filepath}"
    ext = path.suffix.lower()
    extractor = EXTRACTORS.get(ext, extract_txt)
    try:
        text = extractor(str(filepath))
        return text, f"extracted {len(text.split()):,} words from {path.name}"
    except Exception as e:
        return None, str(e)

# ── URL & Cloud handlers ──────────────────────────────────────────────────────

HEADERS = {'User-Agent': 'Mozilla/5.0 (X11; Linux x86_64) AppleWebKit/537.36'}

def resolve_drive_url(url):
    """Convert Google Drive share link to direct download."""
    match = re.search(r'/file/d/([^/]+)', url)
    if match:
        fid = match.group(1)
        return f"https://drive.google.com/uc?export=download&id={fid}"
    match = re.search(r'id=([^&]+)', url)
    if match:
        return f"https://drive.google.com/uc?export=download&id={match.group(1)}"
    return url

def resolve_dropbox_url(url):
    """Convert Dropbox share link to direct download."""
    return url.replace('www.dropbox.com', 'dl.dropboxusercontent.com').replace('?dl=0', '').replace('?dl=1', '')

def resolve_onedrive_url(url):
    """Convert OneDrive share link to direct download."""
    if '1drv.ms' in url or 'onedrive.live.com' in url:
        encoded = urllib.parse.quote(url, safe='')
        return f"https://api.onedrive.com/v1.0/shares/u!{encoded}/root/content"
    return url

def resolve_url(url):
    """Resolve cloud share links to direct download URLs."""
    if 'drive.google.com' in url:
        return resolve_drive_url(url), 'gdrive'
    elif 'dropbox.com' in url:
        return resolve_dropbox_url(url), 'dropbox'
    elif '1drv.ms' in url or 'onedrive.live.com' in url:
        return resolve_onedrive_url(url), 'onedrive'
    return url, 'direct'

def fetch_url(url, max_chars=100000):
    """Fetch text content from any URL. Smart extraction for known sources."""
    resolved_url, source_type = resolve_url(url)

    try:
        r = requests.get(resolved_url, headers=HEADERS, timeout=30, allow_redirects=True)
        if r.status_code != 200:
            return None, f"HTTP {r.status_code}"

        content_type = r.headers.get('content-type', '').lower()

        # Detect file type from content-type or URL
        if 'pdf' in content_type or url.lower().endswith('.pdf'):
            import tempfile
            with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as tmp:
                tmp.write(r.content)
                tmp_path = tmp.name
            text, status = extract_from_file(tmp_path)
            os.unlink(tmp_path)
            return text, status

        elif 'word' in content_type or 'docx' in content_type or url.lower().endswith('.docx'):
            import tempfile
            with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
                tmp.write(r.content)
                tmp_path = tmp.name
            text, status = extract_from_file(tmp_path)
            os.unlink(tmp_path)
            return text, status

        elif 'spreadsheet' in content_type or url.lower().endswith('.xlsx'):
            import tempfile
            with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
                tmp.write(r.content)
                tmp_path = tmp.name
            text, status = extract_from_file(tmp_path)
            os.unlink(tmp_path)
            return text, status

        elif 'html' in content_type:
            # Smart extraction for known high-value sources
            text = _smart_extract_html(url, r.text)
            return text[:max_chars], f"extracted {len(text.split()):,} words from webpage"

        elif 'json' in content_type:
            return r.text[:max_chars], f"extracted JSON ({len(r.text.split()):,} words)"

        else:
            # Try as plain text (covers Project Gutenberg .txt files)
            text = r.text[:max_chars]
            return text, f"extracted {len(text.split()):,} words"

    except Exception as e:
        return None, str(e)


def _smart_extract_html(url, html_text):
    """
    Smart content extraction for known sources.
    Extracts article body text, strips navigation/sidebars/footers.
    Falls back to generic extraction for unknown sources.
    """
    try:
        from bs4 import BeautifulSoup
    except ImportError:
        return extract_html(html_text=html_text)

    soup = BeautifulSoup(html_text, 'html.parser')

    # Wikipedia — extract #mw-content-text article body
    if 'wikipedia.org' in url:
        content = soup.find('div', {'id': 'mw-content-text'})
        if content:
            # Remove navboxes, sidebars, reference lists, edit links
            for tag in content.find_all(['table', 'sup', 'span'], class_=lambda c: c and any(
                x in str(c) for x in ['navbox', 'sidebar', 'reflist', 'mw-editsection', 'noprint', 'metadata'])):
                tag.decompose()
            for tag in content.find_all('div', class_=lambda c: c and any(
                x in str(c) for x in ['reflist', 'navbox', 'sistersitebox', 'noprint'])):
                tag.decompose()
            return content.get_text(separator='\n', strip=True)

    # Stanford Encyclopedia of Philosophy — extract #aueditable article body
    if 'plato.stanford.edu' in url:
        content = soup.find('div', {'id': 'aueditable'}) or soup.find('div', {'id': 'article-content'})
        if content:
            for tag in content.find_all(['nav', 'footer']):
                tag.decompose()
            return content.get_text(separator='\n', strip=True)

    # World History Encyclopedia — extract article body
    if 'worldhistory.org' in url:
        content = soup.find('div', class_='article-content-main') or soup.find('article') or soup.find('div', {'id': 'content'})
        if content:
            for tag in content.find_all(['nav', 'footer', 'aside']):
                tag.decompose()
            return content.get_text(separator='\n', strip=True)

    # arXiv abstract pages — extract abstract text
    if 'arxiv.org' in url:
        abstract = soup.find('blockquote', class_='abstract')
        title = soup.find('h1', class_='title')
        parts = []
        if title:
            parts.append(title.get_text(strip=True))
        if abstract:
            parts.append(abstract.get_text(strip=True))
        if parts:
            return '\n\n'.join(parts)

    # Project Gutenberg HTML — extract body text
    if 'gutenberg.org' in url:
        content = soup.find('div', {'id': 'pg-machine-header'})
        if content:
            content.decompose()
        body = soup.find('body')
        if body:
            for tag in body.find_all(['table', 'pre']):
                if 'START OF' in tag.get_text() or 'END OF' in tag.get_text():
                    tag.decompose()
            return body.get_text(separator='\n', strip=True)

    # Generic fallback — strip scripts/styles/nav/footer
    for tag in soup(['script', 'style', 'nav', 'footer', 'aside', 'header']):
        tag.decompose()
    # Try to find main content area
    main = soup.find('main') or soup.find('article') or soup.find('div', {'role': 'main'})
    if main:
        return main.get_text(separator='\n', strip=True)
    return soup.get_text(separator='\n', strip=True)

# ── Main interface ────────────────────────────────────────────────────────────

def ingest_document(source, max_chars=100000):
    """
    Universal document ingestion.
    source can be: local filepath, URL, Google Drive link, Dropbox link, OneDrive link
    Returns (text, status_message)
    """
    source = source.strip()

    # URL or cloud link
    if source.startswith('http://') or source.startswith('https://'):
        return fetch_url(source, max_chars)

    # Local file
    if os.path.exists(source):
        return extract_from_file(source)

    # Try as path from common locations
    for base in ['/sdcard/Download/', '/sdcard/', os.path.expanduser('~/'), '']:
        full_path = os.path.join(base, source)
        if os.path.exists(full_path):
            return extract_from_file(full_path)

    return None, f"Cannot find: {source}"


def handle_ingest_command(cmd_parts, engine=None, memory=None):
    """
    Handle /ingest commands from chat.py

    /ingest URL                 — fetch and train from URL
    /ingest /path/to/file       — train from local file
    /ingest /sdcard/Download/x  — train from phone storage
    """
    if len(cmd_parts) < 2:
        return """  Ingest Commands:
  /ingest URL              — fetch webpage or file from URL
  /ingest FILEPATH         — ingest local file (any format)
  /ingest /sdcard/Download/FILE — ingest from phone Downloads

  Supported: .txt .pdf .docx .xlsx .pptx .csv .html .json .md .py .png .jpg
  Cloud:     Google Drive, Dropbox, OneDrive share links"""

    source = ' '.join(cmd_parts[1:])
    print(f"  Ingesting: {source[:60]}...")

    text, status = ingest_document(source)

    if not text:
        return f"  Failed: {status}"

    if engine:
        engine.learn_from_text(text)
        # Update memory
        if memory:
            from quantum_language_engine import ConceptExtractor
            extractor = ConceptExtractor()
            extractor.update_corpus_stats(text)
            concepts = extractor.extract_concepts(text, max_concepts=10)
            for c in concepts:
                name = c if isinstance(c, str) else c.get('concept', '')
                if name and name not in memory.concepts:
                    memory.concepts[name] = {'count': 1, 'strength': 1.0}
                    memory.growth['total_concepts'] += 1
            memory.growth['total_insights'] += 1

        word_count = len(text.split())
        return f"  ✓ {status}\n  Trained Markov chain: +{word_count:,} words\n  Chain states: {len(engine.markov.chain):,}"

    return f"  ✓ {status}\n  (no engine — text extracted only)"


if __name__ == "__main__":
    import sys
    if len(sys.argv) > 1:
        source = ' '.join(sys.argv[1:])
        text, status = ingest_document(source)
        print(f"Status: {status}")
        if text:
            print(f"Preview: {text[:500]}")
    else:
        print("Usage: python document_ingester.py URL_OR_FILEPATH")
